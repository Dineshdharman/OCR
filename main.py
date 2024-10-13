import cv2
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import os
import pptx
import openpyxl
import xlrd

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def preprocess_image(image, blur_kernel_size=(7, 7), morph_kernel_size=(3, 50)):
    try:
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        denoised = cv2.fastNlMeansDenoising(gray, None, h=10, searchWindowSize=21, templateWindowSize=7)
        thresh = cv2.adaptiveThreshold(denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY_INV, 11, 2)
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, morph_kernel_size)
        dilate = cv2.dilate(thresh, kernel, iterations=1)
        return dilate
    except Exception as e:
        print(f"Error during image preprocessing: {e}")
        return None

def process_image(image, preprocess_params):
    try:
        base_image = image.copy()
        dilate = preprocess_image(image, *preprocess_params)

        if dilate is not None:
            cnts = cv2.findContours(dilate, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            cnts = cnts[0] if len(cnts) == 2 else cnts[1]
            cnts = sorted(cnts, key=lambda x: cv2.boundingRect(x)[1])
            for c in cnts:
                x,y,w,h = cv2.boundingRect(c)
                if h > 200 and w > 250:
                    roi = base_image[y:y+h, x:x+w]
                    cv2.rectangle(image, (x,y), (x+w, y+h), (36, 255, 12), 2)

            cv2.imwrite("output/processed_image.png", image)
            ocr_result_original = pytesseract.image_to_string(base_image, config=r'--oem 3 --psm 6')
            print(ocr_result_original)
    except Exception as e:
        print(f"Error during image processing: {e}")


def process_pdf(file_path, preprocess_params):
    try:
        # Create 'temp' directory if it doesn't exist
        if not os.path.exists('temp'):
            os.makedirs('temp')
        
        images = convert_from_path(file_path)
        for i, page in enumerate(images):
            page.save(f"temp/page_{i}.jpg", "JPEG")
            image = cv2.imread(f"temp/page_{i}.jpg")
            process_image(image, preprocess_params)
    except Exception as e:
        print(f"Error processing PDF: {e}")


def process_image_file(file_path, preprocess_params):
    try:
        image = cv2.imread(file_path)
        process_image(image, preprocess_params)
    except Exception as e:
        print(f"Error processing image: {e}")

def process_pptx(file_path, preprocess_params):
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            print("Slide content:", " ".join(text))
    except Exception as e:
        print(f"Error processing PowerPoint file: {e}")

def process_xlsx(file_path):
    try:
        workbook = openpyxl.load_workbook(file_path)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        print(cell.value, end=" ")
                print()
    except Exception as e:
        print(f"Error processing Excel file (XLSX): {e}")

def process_xls(file_path):
    try:
        workbook = xlrd.open_workbook(file_path)
        for sheet in workbook.sheets():
            for row_idx in range(sheet.nrows):
                for col_idx in range(sheet.ncols):
                    cell_value = sheet.cell_value(row_idx, col_idx)
                    if cell_value:
                        print(cell_value, end=" ")
                print()
    except Exception as e:
        print(f"Error processing Excel file (XLS): {e}")

def process_bmp(file_path, preprocess_params):
    try:
        image = cv2.imread(file_path)
        process_image(image, preprocess_params)
    except Exception as e:
        print(f"Error processing BMP image: {e}")

def main():
    try:
        file_path = input("Enter the path of file: ")
        file_extension = file_path.split(".")[-1].lower()  # Get file extension in lowercase
        print(f"File extension extracted: {file_extension}")  # Add this line to verify

        # Check for supported formats
        if file_extension == "pdf":
            process_pdf(file_path, preprocess_params=((7, 7), (3, 50)))
        elif file_extension in ["jpg", "jpeg", "png", "tiff", "tif", "bmp"]:  # Include PNG here
            process_image_file(file_path, preprocess_params=((7, 7), (3, 50)))
        elif file_extension == "pptx":
            process_pptx(file_path, preprocess_params=((7, 7), (3, 50)))
        elif file_extension == "xlsx":
            process_xlsx(file_path)
        elif file_extension == "xls":
            process_xls(file_path)
        else:
            print("Unsupported file format. Please provide a PDF, JPG, PNG, TIFF, BMP, PPTX, XLSX, or XLS file.")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    main()